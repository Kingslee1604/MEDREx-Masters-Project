"""
MEDREx Presentation Generator - Clean Light Theme
Run: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors - clean light palette
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
DARK_TEXT  = RGBColor(0x1e, 0x29, 0x3b)
MED_GRAY   = RGBColor(0x47, 0x55, 0x69)
LIGHT_GRAY = RGBColor(0xf1, 0xf5, 0xf9)
BORDER     = RGBColor(0xcb, 0xd5, 0xe1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x16, 0xa3, 0x4a)
GREEN_BG   = RGBColor(0xdc, 0xfc, 0xe7)
ORANGE     = RGBColor(0xea, 0x58, 0x0c)
ORANGE_BG  = RGBColor(0xff, 0xed, 0xd5)
RED        = RGBColor(0xdc, 0x26, 0x26)
RED_BG     = RGBColor(0xfe, 0xe2, 0xe2)
BLUE_BG    = RGBColor(0xef, 0xf6, 0xff)
NAVY_BG    = RGBColor(0xe0, 0xe7, 0xf1)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 11


# ── HELPERS ───────────────────────────────────────────────────────────────────

def bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def bg_color(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def label(slide, text, l, t, w, h, size=14, bold=False,
          color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def multiline(slide, lines, l, t, w, h, default_size=13, default_color=DARK_TEXT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, cfg) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", PP_ALIGN.LEFT)
        p.space_after  = Pt(cfg.get("sa", 3))
        p.space_before = Pt(cfg.get("sb", 0))
        r = p.add_run()
        r.text = text
        r.font.size    = Pt(cfg.get("size", default_size))
        r.font.bold    = cfg.get("bold", False)
        r.font.italic  = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", default_color)

def top_bar(slide, title, subtitle=None):
    """Navy top bar with white title."""
    rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), fill=NAVY)
    label(slide, title,
          Inches(0.4), Inches(0.12), Inches(10), Inches(0.6),
          size=26, bold=True, color=WHITE)
    if subtitle:
        label(slide, subtitle,
              Inches(0.4), Inches(0.72), Inches(10), Inches(0.32),
              size=13, color=RGBColor(0xa8,0xc8,0xe8))

def slide_num(slide, n):
    label(slide, f"{n} / {TOTAL}",
          Inches(12.3), Inches(7.2), Inches(0.9), Inches(0.25),
          size=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)

def stat_box(slide, l, t, w, h, value, caption, v_color=NAVY, fill=LIGHT_GRAY, border=BORDER):
    rect(slide, l, t, w, h, fill=fill, border=border)
    label(slide, value, l, t + Inches(0.18), w, Inches(0.7),
          size=34, bold=True, color=v_color, align=PP_ALIGN.CENTER)
    label(slide, caption, l, t + Inches(0.85), w, Inches(0.4),
          size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

def row_item(slide, t, num, title, desc, status, status_color, status_bg):
    """Single method row."""
    rect(slide, Inches(0.4), t, Inches(12.5), Inches(0.82), fill=WHITE, border=BORDER, bw=Pt(0.75))
    # number circle
    rect(slide, Inches(0.55), t + Inches(0.16), Inches(0.5), Inches(0.5), fill=NAVY_BG)
    label(slide, str(num), Inches(0.55), t + Inches(0.2), Inches(0.5), Inches(0.4),
          size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # title
    label(slide, title, Inches(1.2), t + Inches(0.1), Inches(6.2), Inches(0.38),
          size=14, bold=True, color=DARK_TEXT)
    label(slide, desc, Inches(1.2), t + Inches(0.48), Inches(6.2), Inches(0.28),
          size=11, color=MED_GRAY)
    # status badge
    rect(slide, Inches(7.6), t + Inches(0.2), Inches(1.8), Inches(0.42),
         fill=status_bg, border=status_color, bw=Pt(1))
    label(slide, status, Inches(7.6), t + Inches(0.23), Inches(1.8), Inches(0.35),
          size=12, bold=True, color=status_color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)

# Left navy accent strip
rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), fill=NAVY)

# Top light bar
rect(s, Inches(0.08), Inches(0), Inches(13.25), Inches(0.06), fill=NAVY)

# Main title area
label(s, "MEDREx",
      Inches(0.5), Inches(0.8), Inches(9), Inches(1.5),
      size=72, bold=True, color=NAVY)

label(s, "Medical Evidence & Decision Reasoning eXchange",
      Inches(0.5), Inches(2.25), Inches(9), Inches(0.55),
      size=20, bold=False, color=DARK_TEXT)

label(s, "LLM-Based Diagnostic Insight Extraction from Pathology Reports",
      Inches(0.5), Inches(2.85), Inches(9), Inches(0.45),
      size=15, italic=True, color=MED_GRAY)

# Divider line
rect(s, Inches(0.5), Inches(3.45), Inches(5.5), Pt(2), fill=NAVY)

# Info block
info = [
    ("Student",  "Kingslee Dominic Savio Velu"),
    ("Course",   "CSC 590 – Master's Project, CSUDH Spring 2026"),
    ("Mentor",   "Dr. Bin Tang")
]
for i, (key, val) in enumerate(info):
    top = Inches(3.65 + i * 0.58)
    label(s, key + ":",  Inches(0.5), top, Inches(1.2), Inches(0.45),
          size=13, bold=True, color=NAVY)
    label(s, val,        Inches(1.7), top, Inches(7.5), Inches(0.45),
          size=13, color=DARK_TEXT)

# Right side stat cards
for i, (val, cap, vc) in enumerate([
    ("9,523",   "Pathology Reports", NAVY),
    ("35",      "Cancer Types",       NAVY),
    ("96.36%",  "TF-IDF Accuracy",    GREEN),
]):
    stat_box(s, Inches(10.0), Inches(1.3 + i * 1.85), Inches(3.0), Inches(1.55),
             val, cap, v_color=vc)

slide_num(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "The Problem", "Why automated classification of pathology reports is hard")
slide_num(s, 2)

problems = [
    ("Pathology reports are unstructured, long-form text — difficult to classify manually at scale",
     ORANGE, ORANGE_BG),
    ("35 cancer types with overlapping medical terminology — e.g. COAD vs READ both say 'adenocarcinoma'",
     ORANGE, ORANGE_BG),
    ("Patient data cannot be sent to cloud APIs (ChatGPT, etc.) — HIPAA violation",
     RED, RED_BG),
    ("LLMs hallucinate — they give confident but wrong answers with no source citations",
     RED, RED_BG),
    ("Need: explainable, evidence-grounded, locally-run classification system",
     GREEN, GREEN_BG),
]
for i, (text, col, bg_c) in enumerate(problems):
    t = Inches(1.3 + i * 1.05)
    rect(s, Inches(0.4), t, Inches(12.5), Inches(0.82), fill=bg_c, border=col, bw=Pt(1))
    label(s, text, Inches(0.65), t + Inches(0.2), Inches(12.0), Inches(0.45),
          size=15, color=DARK_TEXT)

rect(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.65), fill=NAVY)
label(s, "MEDREx solves this by comparing 5 NLP methods scientifically and building a RAG pipeline that is accurate, local, and explainable.",
      Inches(0.6), Inches(6.7), Inches(12.0), Inches(0.45),
      size=13, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "Dataset — TCGA Pathology Reports",
        "The Cancer Genome Atlas | National Cancer Institute")
slide_num(s, 3)

# Stat cards
stats = [
    ("9,523",  "Total Reports",          NAVY,  BLUE_BG,  NAVY),
    ("35",     "Cancer Types",           NAVY,  BLUE_BG,  NAVY),
    ("560",    "Avg Words / Report",     NAVY,  LIGHT_GRAY, BORDER),
    ("42.7%",  "Reports > 512 tokens",  ORANGE, ORANGE_BG, ORANGE),
    ("24x",    "Class Imbalance Ratio",  ORANGE, ORANGE_BG, ORANGE),
    ("96.36%", "Best Accuracy (TF-IDF)", GREEN,  GREEN_BG,  GREEN),
]
for i, (val, cap, vc, bg_c, bdr) in enumerate(stats):
    col_i = i % 3; row_i = i // 3
    stat_box(s,
             Inches(0.4 + col_i * 4.22), Inches(1.3 + row_i * 2.1),
             Inches(3.9), Inches(1.75),
             val, cap, v_color=vc, fill=bg_c, border=bdr)

# Bottom note
rect(s, Inches(0.4), Inches(5.62), Inches(12.5), Pt(1.5), fill=BORDER)
multiline(s, [
    ("Dataset: TCGA_Reports.csv (34MB)   |   Labels: tcga_patient_to_cancer_type.csv   |   Joined on patient_id",
     {"size": 12, "color": MED_GRAY, "align": PP_ALIGN.CENTER}),
    ("Largest class: BRCA — 1,034 reports   |   Smallest: CHOL — 43 reports   |   Handled with class_weight='balanced'",
     {"size": 12, "color": MED_GRAY, "align": PP_ALIGN.CENTER, "sb": 4}),
],
Inches(0.5), Inches(5.75), Inches(12.2), Inches(0.7))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "System Architecture", "FastAPI backend + Streamlit UI + ChromaDB vector store")
slide_num(s, 4)

# Three columns
cols = [
    ("1  Data Ingestion", NAVY, BLUE_BG,
     ["9,523 TCGA reports loaded", "Sentence-transformer model",
      "all-MiniLM-L6-v2 (384-dim)", "Embed & store in ChromaDB",
      "Runs once, saved to disk"]),
    ("2  Retrieval", GREEN, GREEN_BG,
     ["ChromaDB vector search", "Cosine similarity query",
      "Top-k similar reports", "Returns metadata + text",
      "9,523 reports indexed"]),
    ("3  Generation  (Planned)", ORANGE, ORANGE_BG,
     ["LLM receives query + cases", "Generates grounded answer",
      "Source citations included", "Runs locally — HIPAA safe",
      "LLaMA / BioMedLM"]),
]
for i, (title, col, bg_c, items) in enumerate(cols):
    l = Inches(0.4 + i * 4.22)
    rect(s, l, Inches(1.3), Inches(4.0), Inches(4.5), fill=bg_c, border=col, bw=Pt(1.5))
    label(s, title, l + Inches(0.15), Inches(1.38), Inches(3.7), Inches(0.5),
          size=16, bold=True, color=col)
    rect(s, l + Inches(0.1), Inches(1.92), Inches(3.8), Pt(1.5), fill=col)
    lines = [("• " + item, {"size": 13, "color": DARK_TEXT, "sa": 8})
             for item in items]
    multiline(s, lines, l + Inches(0.2), Inches(2.05), Inches(3.6), Inches(3.5))

# Arrows
for i in range(2):
    label(s, "->", Inches(4.37 + i * 4.22), Inches(3.25), Inches(0.3), Inches(0.5),
          size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Bottom bar
rect(s, Inches(0.4), Inches(6.05), Inches(12.5), Pt(1), fill=BORDER)
label(s, "Tech Stack:  Python  |  FastAPI  |  Streamlit  |  ChromaDB  |  sentence-transformers  |  scikit-learn  |  LangChain (planned)",
      Inches(0.5), Inches(6.18), Inches(12.2), Inches(0.38),
      size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.6), fill=RED_BG, border=RED, bw=Pt(1))
label(s, "Why not ChatGPT?  HIPAA prohibits sending patient data to external APIs. Also: academic goal is to scientifically compare methods.",
      Inches(0.6), Inches(6.73), Inches(12.0), Inches(0.38),
      size=12, color=RED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW ACCURACY IS CALCULATED
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "How Accuracy is Calculated", "Train/Test Split -> Predict -> Score")
slide_num(s, 5)

steps = [
    ("Load Data",    "9,523 reports\nwith labels"),
    ("Split 70/30",  "Stratified split\n70% train, 30% test"),
    ("Train Model",  "Model learns from\ntraining data only"),
    ("Predict",      "Predict cancer type\nfor test reports"),
    ("Compare",      "Prediction vs\ntrue label"),
    ("Score",        "Correct / Total\n= Accuracy %"),
]
for i, (title, desc) in enumerate(steps):
    l = Inches(0.35 + i * 2.1)
    fill = GREEN_BG if i == 5 else LIGHT_GRAY
    bdr  = GREEN    if i == 5 else BORDER
    rect(s, l, Inches(1.3), Inches(1.9), Inches(2.1), fill=fill, border=bdr, bw=Pt(1))
    label(s, title, l, Inches(1.42), Inches(1.9), Inches(0.45),
          size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    label(s, desc, l, Inches(1.9), Inches(1.9), Inches(0.9),
          size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        label(s, "->", Inches(2.22 + i * 2.1), Inches(2.1),
              Inches(0.22), Inches(0.4),
              size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Divider
rect(s, Inches(0.4), Inches(3.65), Inches(12.5), Pt(1.5), fill=BORDER)

# Three metric boxes
metrics = [
    ("Accuracy",
     "Correct predictions divided by total predictions.\nOur result: 96.36%",
     GREEN, GREEN_BG),
    ("Weighted F1",
     "Balances precision and recall across all 35 classes.\nHandles class imbalance fairly.",
     NAVY, BLUE_BG),
    ("Per-Class F1",
     "F1 score for each individual cancer type.\nShows which types the model struggles with.",
     ORANGE, ORANGE_BG),
]
for i, (title, desc, col, bg_c) in enumerate(metrics):
    l = Inches(0.4 + i * 4.22)
    rect(s, l, Inches(3.85), Inches(4.0), Inches(2.2), fill=bg_c, border=col, bw=Pt(1))
    label(s, title, l + Inches(0.15), Inches(3.95), Inches(3.7), Inches(0.45),
          size=17, bold=True, color=col)
    label(s, desc, l + Inches(0.15), Inches(4.45), Inches(3.7), Inches(1.3),
          size=13, color=DARK_TEXT)

label(s, "Stratified split: All 35 cancer types appear in both train and test sets — no class is excluded from evaluation.",
      Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.4),
      size=12, italic=True, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TF-IDF RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "Method 1 — TF-IDF + Logistic Regression",
        "COMPLETED  |  Beats Cedars-Sinai reference result")
slide_num(s, 6)

# Comparison boxes
rect(s, Inches(0.4), Inches(1.3), Inches(5.9), Inches(3.1), fill=GREEN_BG, border=GREEN, bw=Pt(2))
label(s, "MEDREx Result", Inches(0.5), Inches(1.42), Inches(5.7), Inches(0.4),
      size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)
label(s, "96.36%", Inches(0.5), Inches(1.82), Inches(5.7), Inches(1.2),
      size=60, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
label(s, "TF-IDF + Logistic Regression\n35 cancer types  |  30% test set",
      Inches(0.5), Inches(3.05), Inches(5.7), Inches(0.7),
      size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

rect(s, Inches(6.7), Inches(1.3), Inches(5.9), Inches(3.1), fill=BLUE_BG, border=NAVY, bw=Pt(1.5))
label(s, "Cedars-Sinai Reference", Inches(6.8), Inches(1.42), Inches(5.7), Inches(0.4),
      size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)
label(s, "95.31%", Inches(6.8), Inches(1.82), Inches(5.7), Inches(1.2),
      size=60, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
label(s, "BoW + Logistic Regression\n32 cancer types  |  reference benchmark",
      Inches(6.8), Inches(3.05), Inches(5.7), Inches(0.7),
      size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(4.55), Inches(12.2), Inches(0.55), fill=GREEN_BG, border=GREEN, bw=Pt(1))
label(s, "+1.05% improvement   |   3 more cancer types covered   |   Better TF-IDF features",
      Inches(0.5), Inches(4.65), Inches(12.0), Inches(0.35),
      size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Config grid
config = [
    ("max_features=15000",    "15K vocabulary features"),
    ("ngram_range=(1,2)",     "Unigrams + Bigrams"),
    ("sublinear_tf=True",     "Log scaling on term freq"),
    ("class_weight='balanced'","Handles 24x imbalance"),
    ("C=1.0, max_iter=1000",  "LR solver: lbfgs"),
    ("Saved to disk",         "joblib — instant reload"),
]
for i, (code, desc) in enumerate(config):
    col_i = i % 3; row_i = i // 3
    l = Inches(0.4 + col_i * 4.12)
    t = Inches(5.3 + row_i * 0.65)
    rect(s, l, t, Inches(3.9), Inches(0.55), fill=LIGHT_GRAY, border=BORDER, bw=Pt(0.75))
    label(s, code, l + Inches(0.12), t + Inches(0.1), Inches(2.0), Inches(0.35),
          size=12, bold=True, color=NAVY)
    label(s, desc, l + Inches(2.15), t + Inches(0.1), Inches(1.65), Inches(0.35),
          size=11, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CHROMADB
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "ChromaDB Vector Store — COMPLETED",
        "All 9,523 reports embedded and indexed with cosine similarity")
slide_num(s, 7)

# Status banner
rect(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.6), fill=GREEN_BG, border=GREEN, bw=Pt(1.5))
label(s, "9,523 / 9,523 reports embedded   |   vector_ready: true   |   Saved to disk at backend/vector_store/",
      Inches(0.6), Inches(1.42), Inches(12.0), Inches(0.38),
      size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# How it works — pipeline
steps = [
    ("Report Text",    "Input pathology\nreport (raw text)"),
    ("Truncate",       "First 1000 chars\n(consistent length)"),
    ("Encode",         "all-MiniLM-L6-v2\n-> 384-dim vector"),
    ("Store",          "ChromaDB saves\nvector + metadata"),
    ("Query",          "New text encoded\n-> cosine search"),
]
for i, (title, desc) in enumerate(steps):
    l = Inches(0.4 + i * 2.46)
    rect(s, l, Inches(2.1), Inches(2.2), Inches(1.8), fill=LIGHT_GRAY, border=BORDER, bw=Pt(1))
    label(s, title, l, Inches(2.2), Inches(2.2), Inches(0.42),
          size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    label(s, desc, l, Inches(2.65), Inches(2.2), Inches(0.9),
          size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        label(s, "->", Inches(2.55 + i * 2.46), Inches(2.8),
              Inches(0.25), Inches(0.38),
              size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Three info boxes
infos = [
    ("Why ChromaDB?",
     ["Persists to disk — no re-embedding on restart",
      "Runs 100% locally — fully HIPAA compliant",
      "Built-in cosine similarity search",
      "Stores metadata: cancer_type, word_count"],
     NAVY, BLUE_BG),
    ("What it enables",
     ["Method 4: Embedding Similarity (k-NN vote)",
      "Method 5: RAG — retrieve cases for LLM",
      "Semantic search across all 9,523 reports",
      "Similar cases shown in Live Prediction tab"],
     GREEN, GREEN_BG),
    ("Embedding Model",
     ["Model: all-MiniLM-L6-v2 (HuggingFace)",
      "Size: ~80MB, runs on CPU",
      "Output: 384-dimensional float vector",
      "Medical alt: BioBERT (available if needed)"],
     ORANGE, ORANGE_BG),
]
for i, (title, items, col, bg_c) in enumerate(infos):
    l = Inches(0.4 + i * 4.22)
    rect(s, l, Inches(4.15), Inches(4.0), Inches(2.8), fill=bg_c, border=col, bw=Pt(1))
    label(s, title, l + Inches(0.15), Inches(4.25), Inches(3.7), Inches(0.42),
          size=15, bold=True, color=col)
    rect(s, l + Inches(0.1), Inches(4.72), Inches(3.8), Pt(1.5), fill=col)
    lines = [("• " + item, {"size": 12, "color": DARK_TEXT, "sa": 7})
             for item in items]
    multiline(s, lines, l + Inches(0.15), Inches(4.85), Inches(3.7), Inches(1.9))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PER-CLASS F1
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "Per-Class F1 — TF-IDF + LR Results",
        "Which cancer types are hardest to classify, and why?")
slide_num(s, 8)

# Hardest
label(s, "Hardest to Predict", Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.4),
      size=15, bold=True, color=RED)
hard = [
    ("UCS",  "0.80", "Uterine Carcinosarcoma",
     "Language overlaps with UCEC — both uterine cancers"),
    ("READ", "0.84", "Rectum Adenocarcinoma",
     "Nearly identical reports to COAD (Colon)"),
    ("LUSC", "0.91", "Lung Squamous Cell Carcinoma",
     "Overlaps with LUAD — both lung, different cell type"),
]
for i, (code, score, name, reason) in enumerate(hard):
    t = Inches(1.8 + i * 1.1)
    col = RED if float(score) < 0.85 else ORANGE
    bg_c = RED_BG if float(score) < 0.85 else ORANGE_BG
    rect(s, Inches(0.4), t, Inches(6.0), Inches(0.88), fill=bg_c, border=col, bw=Pt(1))
    label(s, code, Inches(0.55), t + Inches(0.22), Inches(0.7), Inches(0.4),
          size=15, bold=True, color=col)
    label(s, score, Inches(1.35), t + Inches(0.22), Inches(0.7), Inches(0.4),
          size=15, bold=True, color=col)
    label(s, f"{name} — {reason}",
          Inches(2.15), t + Inches(0.22), Inches(4.1), Inches(0.4),
          size=12, color=DARK_TEXT)

# Perfect F1
label(s, "Perfect F1 = 1.00", Inches(0.4), Inches(5.15), Inches(5.5), Inches(0.4),
      size=15, bold=True, color=GREEN)
perfect = ["BRCA", "DLBC", "MESO", "PRAD", "TGCT", "UVM", "LAML", "PCPG", "BLCA", "THCA"]
for i, code in enumerate(perfect):
    col_i = i % 5; row_i = i // 5
    l = Inches(0.4 + col_i * 1.18)
    t = Inches(5.62 + row_i * 0.6)
    rect(s, l, t, Inches(1.05), Inches(0.48), fill=GREEN_BG, border=GREEN, bw=Pt(1))
    label(s, code, l, t + Inches(0.08), Inches(1.05), Inches(0.35),
          size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Right — insight box
rect(s, Inches(6.7), Inches(1.3), Inches(6.2), Inches(5.9),
     fill=LIGHT_GRAY, border=BORDER, bw=Pt(1))
label(s, "Key Insights", Inches(6.85), Inches(1.42), Inches(5.9), Inches(0.42),
      size=16, bold=True, color=NAVY)
rect(s, Inches(6.8), Inches(1.9), Inches(5.9), Pt(1.5), fill=NAVY)
insights = [
    ("TF-IDF uses word frequency — it misses semantic meaning between similar cancer types.", {}),
    ("", {}),
    ("UCS and READ fail because their reports use nearly identical words to a neighboring cancer type.", {}),
    ("", {}),
    ("BERT understands full sentence context — expected to fix these hard cases.", {}),
    ("", {}),
    ("RAG retrieves the most similar real cases — if those cases are correct, the LLM prediction improves.", {}),
    ("", {}),
    ("This per-class analysis scientifically justifies why we go beyond TF-IDF to BERT and RAG.", {}),
]
multiline(s, insights, Inches(6.85), Inches(2.05), Inches(5.7), Inches(4.8),
          default_size=13, default_color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 5-METHOD ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "5-Method Comparison Roadmap",
        "Scientific framework to justify the MEDREx RAG architecture")
slide_num(s, 9)

# Column headers
for lbl, l in [("Method", Inches(1.0)), ("Category", Inches(6.0)),
               ("Status", Inches(7.65)), ("Accuracy", Inches(9.5)),
               ("Notes", Inches(11.1))]:
    label(s, lbl, l, Inches(1.28), Inches(1.7), Inches(0.3),
          size=11, bold=True, color=MED_GRAY)
rect(s, Inches(0.4), Inches(1.6), Inches(12.5), Pt(1.5), fill=BORDER)

methods = [
    (1, "TF-IDF + Logistic Regression", "Traditional ML",
     "DONE", GREEN, GREEN_BG,
     "96.36%", GREEN,
     "Fast, explainable, no GPU needed"),
    (2, "TF-IDF + SVM",                 "Traditional ML",
     "NEXT", ORANGE, ORANGE_BG,
     "TBD", MED_GRAY,
     "Often outperforms LR on text"),
    (3, "BERT (Bio_ClinicalBERT)",       "Transformer",
     "PLANNED", NAVY, BLUE_BG,
     "TBD", MED_GRAY,
     "Understands medical context"),
    (4, "Embedding Similarity (k-NN)",   "Semantic Search",
     "READY", GREEN, GREEN_BG,
     "TBD", MED_GRAY,
     "ChromaDB fully built"),
    (5, "RAG + LLM  (MEDREx Core)",       "LLM + Retrieval",
     "PLANNED", NAVY, BLUE_BG,
     "TBD", MED_GRAY,
     "Full MEDREx pipeline"),
]
for i, (num, name, cat, status, sc, sbg, acc, ac, note) in enumerate(methods):
    t = Inches(1.72 + i * 1.05)
    rect(s, Inches(0.4), t, Inches(12.5), Inches(0.88),
         fill=WHITE, border=BORDER, bw=Pt(0.75))
    # num
    rect(s, Inches(0.55), t + Inches(0.17), Inches(0.45), Inches(0.52), fill=NAVY_BG)
    label(s, str(num), Inches(0.55), t + Inches(0.22), Inches(0.45), Inches(0.38),
          size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # name
    label(s, name, Inches(1.1), t + Inches(0.12), Inches(4.8), Inches(0.38),
          size=14, bold=True, color=DARK_TEXT)
    label(s, cat, Inches(6.05), t + Inches(0.25), Inches(1.5), Inches(0.35),
          size=12, color=MED_GRAY)
    # status badge
    rect(s, Inches(7.65), t + Inches(0.2), Inches(1.7), Inches(0.42),
         fill=sbg, border=sc, bw=Pt(1))
    label(s, status, Inches(7.65), t + Inches(0.24), Inches(1.7), Inches(0.35),
          size=12, bold=True, color=sc, align=PP_ALIGN.CENTER)
    # accuracy
    label(s, acc, Inches(9.5), t + Inches(0.22), Inches(1.5), Inches(0.38),
          size=15, bold=True, color=ac, align=PP_ALIGN.CENTER)
    # note
    label(s, note, Inches(11.1), t + Inches(0.25), Inches(1.7), Inches(0.35),
          size=11, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — UPCOMING PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "Upcoming Plan", "What's next after today's demo")
slide_num(s, 10)

phases = [
    ("Phase 2 — TF-IDF + SVM", ORANGE, ORANGE_BG,
     ["Add LinearSVC as second method",
      "Compare directly vs LR on same test set",
      "Expected accuracy: 94–97%",
      "No GPU needed — runs fast"]),
    ("Phase 3 — BERT Fine-tuning", NAVY, BLUE_BG,
     ["Use Bio_ClinicalBERT or ClinicalBigBird",
      "Fine-tune for 35-class cancer classification",
      "Needs GPU (Google Colab / local)",
      "Expected to fix UCS, READ hard cases"]),
    ("Phase 4 — Embedding Similarity", GREEN, GREEN_BG,
     ["ChromaDB is 100% built and ready",
      "Run k-NN classification experiment",
      "Majority vote from top-5 similar cases",
      "Compare accuracy vs TF-IDF"]),
]
for i, (title, col, bg_c, items) in enumerate(phases):
    l = Inches(0.4 + i * 4.22)
    rect(s, l, Inches(1.3), Inches(4.0), Inches(4.3), fill=bg_c, border=col, bw=Pt(1.5))
    label(s, title, l + Inches(0.15), Inches(1.42), Inches(3.7), Inches(0.45),
          size=15, bold=True, color=col)
    rect(s, l + Inches(0.1), Inches(1.92), Inches(3.8), Pt(1.5), fill=col)
    lines = [("• " + item, {"size": 13, "color": DARK_TEXT, "sa": 10})
             for item in items]
    multiline(s, lines, l + Inches(0.15), Inches(2.05), Inches(3.7), Inches(3.4))

# Phase 5 — full width
rect(s, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.45),
     fill=BLUE_BG, border=NAVY, bw=Pt(1.5))
label(s, "Phase 5 — RAG + LLM  (MEDREx Core)",
      Inches(0.6), Inches(5.9), Inches(6.0), Inches(0.42),
      size=16, bold=True, color=NAVY)
label(s, "Retrieve top-k similar cases from ChromaDB  ->  Pass as context to local LLM (LLaMA / BioMedLM)  ->  Generate evidence-grounded answer  ->  Full MEDREx system live",
      Inches(0.6), Inches(6.38), Inches(12.0), Inches(0.65),
      size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "Summary", "What has been built, what results achieved, what is coming next")
slide_num(s, 11)

# Completed
rect(s, Inches(0.4), Inches(1.3), Inches(5.9), Inches(4.8),
     fill=GREEN_BG, border=GREEN, bw=Pt(1.5))
label(s, "Completed", Inches(0.55), Inches(1.42), Inches(5.6), Inches(0.42),
      size=17, bold=True, color=GREEN)
done = [
    "FastAPI backend + Streamlit demo UI",
    "TF-IDF + LR: 96.36% accuracy (beats Cedars-Sinai 95.31%)",
    "ChromaDB: all 9,523 reports fully embedded",
    "Per-class F1 analysis across all 35 cancer types",
    "5-method evaluation framework defined",
    "Live Prediction tab with similar cases from ChromaDB",
]
lines = [("  " + item, {"size": 13, "color": DARK_TEXT, "sa": 9}) for item in done]
multiline(s, lines, Inches(0.55), Inches(1.95), Inches(5.6), Inches(3.8))

# Upcoming
rect(s, Inches(6.7), Inches(1.3), Inches(5.9), Inches(4.8),
     fill=BLUE_BG, border=NAVY, bw=Pt(1.5))
label(s, "Upcoming", Inches(6.85), Inches(1.42), Inches(5.6), Inches(0.42),
      size=17, bold=True, color=NAVY)
coming = [
    "TF-IDF + SVM  (Method 2)",
    "BERT fine-tuning — Bio_ClinicalBERT  (Method 3)",
    "Embedding Similarity experiment  (Method 4)",
    "RAG + LLM full pipeline  (Method 5)",
    "Clinical analysis — survival, error analysis",
    "Masters thesis write-up",
]
lines2 = [("  " + item, {"size": 13, "color": DARK_TEXT, "sa": 9}) for item in coming]
multiline(s, lines2, Inches(6.85), Inches(1.95), Inches(5.6), Inches(3.8))

# Closing quote
rect(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.9), fill=LIGHT_GRAY, border=BORDER, bw=Pt(1))
label(s,
      '"MEDREx demonstrates that combining classical NLP with semantic search and LLM reasoning '
      'produces a HIPAA-safe, explainable, and accurate clinical decision support system."',
      Inches(0.6), Inches(6.42), Inches(12.0), Inches(0.65),
      size=14, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── SAVE ──────────────────────────────────────────────────────────────────────
OUT = (r"c:\Users\kings\OneDrive\Desktop\CSUDH_Kingslee"
       r"\CSUDH-Spring-2026\CSC-590-MastersProject"
       r"\MastersProject\MEDREx_Presentation_v2.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"{TOTAL} slides | Widescreen 16:9 | Light theme")
