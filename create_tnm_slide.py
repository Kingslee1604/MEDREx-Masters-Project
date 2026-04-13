"""
Creates a single slide explaining Task 1 vs Task 2 (TNM Staging)
Output: TNM_Slide.pptx  — open in PowerPoint, copy the slide into main deck
Run: python create_tnm_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
DARK_TEXT  = RGBColor(0x1e, 0x29, 0x3b)
MED_GRAY   = RGBColor(0x47, 0x55, 0x69)
LIGHT_GRAY = RGBColor(0xf1, 0xf5, 0xf9)
BORDER     = RGBColor(0xcb, 0xd5, 0xe1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x16, 0xa3, 0x4a)
GREEN_BG   = RGBColor(0xdc, 0xfc, 0xe7)
ORANGE     = RGBColor(0xea, 0x58, 0x0c)
ORANGE_BG  = RGBColor(0xff, 0xed, 0xd5)
BLUE_BG    = RGBColor(0xef, 0xf6, 0xff)
NAVY_BG    = RGBColor(0xe0, 0xe7, 0xf1)
PURPLE     = RGBColor(0x7c, 0x3a, 0xed)
PURPLE_BG  = RGBColor(0xed, 0xe9, 0xfe)
RED        = RGBColor(0xdc, 0x26, 0x26)
RED_BG     = RGBColor(0xfe, 0xe2, 0xe2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def label(slide, text, l, t, w, h, size=13, bold=False,
          color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def multiline(slide, lines, l, t, w, h, default_size=13, default_color=DARK_TEXT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, cfg) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", PP_ALIGN.LEFT)
        p.space_after  = Pt(cfg.get("sa", 3))
        p.space_before = Pt(cfg.get("sb", 0))
        r = p.add_run()
        r.text = text
        r.font.size    = Pt(cfg.get("size", default_size))
        r.font.bold    = cfg.get("bold", False)
        r.font.italic  = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", default_color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE A — TWO TASKS OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)

# Top bar
rect(s, Inches(0), Inches(0), Inches(13.33), Inches(1.1), fill=NAVY)
label(s, "Two Core Tasks — What MEDREx is Solving",
      Inches(0.4), Inches(0.12), Inches(10), Inches(0.6),
      size=26, bold=True, color=WHITE)
label(s, "Both tasks use the same pathology report text as input — but answer different clinical questions",
      Inches(0.4), Inches(0.72), Inches(12), Inches(0.32),
      size=13, color=RGBColor(0xa8,0xc8,0xe8))

# TASK 1 box
rect(s, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.9),
     fill=BLUE_BG, border=NAVY, bw=Pt(2))
label(s, "TASK 1 — Cancer Type Classification",
      Inches(0.55), Inches(1.42), Inches(5.7), Inches(0.45),
      size=16, bold=True, color=NAVY)
rect(s, Inches(0.5), Inches(1.92), Inches(5.8), Pt(1.5), fill=NAVY)

label(s, "Question:",
      Inches(0.55), Inches(2.05), Inches(1.2), Inches(0.35),
      size=13, bold=True, color=NAVY)
label(s, "Which of the 35 cancer types does this report belong to?",
      Inches(1.75), Inches(2.05), Inches(4.5), Inches(0.35),
      size=13, color=DARK_TEXT)

label(s, "Input:",
      Inches(0.55), Inches(2.52), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "Pathology report text",
      Inches(1.55), Inches(2.52), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

label(s, "Output:",
      Inches(0.55), Inches(2.88), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "Cancer type code  e.g. BRCA, LUAD, GBM",
      Inches(1.55), Inches(2.88), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

label(s, "Labels:",
      Inches(0.55), Inches(3.24), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "tcga_patient_to_cancer_type.csv  (11,159 rows)",
      Inches(1.55), Inches(3.24), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

label(s, "Classes:",
      Inches(0.55), Inches(3.6), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "35 cancer types",
      Inches(1.55), Inches(3.6), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

# Status
rect(s, Inches(0.55), Inches(4.05), Inches(5.6), Inches(1.5),
     fill=GREEN_BG, border=GREEN, bw=Pt(1))
label(s, "STATUS — IN PROGRESS",
      Inches(0.7), Inches(4.12), Inches(5.2), Inches(0.35),
      size=12, bold=True, color=GREEN)
multiline(s, [
    ("Method 1 — TF-IDF + LR:  96.36%  DONE", {"size": 12, "color": GREEN, "bold": True, "sa": 5}),
    ("Method 2 — TF-IDF + SVM:  TBD  (Next)", {"size": 12, "color": ORANGE, "sa": 5}),
    ("Method 3 — BERT:  Planned", {"size": 12, "color": MED_GRAY, "sa": 5}),
    ("Method 4 — Embedding Similarity:  ChromaDB Ready", {"size": 12, "color": NAVY, "sa": 5}),
    ("Method 5 — RAG + LLM:  Planned", {"size": 12, "color": MED_GRAY, "sa": 5}),
], Inches(0.7), Inches(4.52), Inches(5.2), Inches(0.95))

# Cedars ref
rect(s, Inches(0.55), Inches(5.65), Inches(5.6), Inches(0.5),
     fill=NAVY_BG, border=NAVY, bw=Pt(1))
label(s, "Cedars-Sinai reference: BoW + LR = 95.31%  (32 classes)  |  MEDREx beats it",
      Inches(0.7), Inches(5.72), Inches(5.3), Inches(0.35),
      size=11, color=NAVY)


# TASK 2 box
rect(s, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.9),
     fill=PURPLE_BG, border=PURPLE, bw=Pt(2))
label(s, "TASK 2 — TNM Stage Classification",
      Inches(7.05), Inches(1.42), Inches(5.7), Inches(0.45),
      size=16, bold=True, color=PURPLE)
rect(s, Inches(7.0), Inches(1.92), Inches(5.8), Pt(1.5), fill=PURPLE)

label(s, "Question:",
      Inches(7.05), Inches(2.05), Inches(1.2), Inches(0.35),
      size=13, bold=True, color=PURPLE)
label(s, "What is the T, N, M severity stage from this report?",
      Inches(8.25), Inches(2.05), Inches(4.5), Inches(0.35),
      size=13, color=DARK_TEXT)

label(s, "Input:",
      Inches(7.05), Inches(2.52), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "Same pathology report text",
      Inches(8.05), Inches(2.52), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

label(s, "Output:",
      Inches(7.05), Inches(2.88), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "3 separate labels: T stage + N stage + M stage",
      Inches(8.05), Inches(2.88), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

label(s, "Labels:",
      Inches(7.05), Inches(3.24), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "3 CSV files in TNM_Stage/ folder",
      Inches(8.05), Inches(3.24), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

label(s, "Classes:",
      Inches(7.05), Inches(3.6), Inches(1.0), Inches(0.32),
      size=12, bold=True, color=MED_GRAY)
label(s, "T: T1-T4  |  N: N0-N3  |  M: M0, M1",
      Inches(8.05), Inches(3.6), Inches(4.6), Inches(0.32),
      size=12, color=DARK_TEXT)

# Status
rect(s, Inches(7.05), Inches(4.05), Inches(5.6), Inches(1.5),
     fill=ORANGE_BG, border=ORANGE, bw=Pt(1))
label(s, "STATUS — NOT STARTED YET",
      Inches(7.2), Inches(4.12), Inches(5.2), Inches(0.35),
      size=12, bold=True, color=ORANGE)
multiline(s, [
    ("T-stage model:  Not built  (6,966 labeled rows)", {"size": 12, "color": MED_GRAY, "sa": 5}),
    ("N-stage model:  Not built  (5,678 labeled rows)", {"size": 12, "color": MED_GRAY, "sa": 5}),
    ("M-stage model:  Not built  (4,608 labeled rows)", {"size": 12, "color": MED_GRAY, "sa": 5}),
    ("Will follow same 5-method comparison as Task 1", {"size": 12, "color": MED_GRAY, "sa": 5}),
], Inches(7.2), Inches(4.52), Inches(5.2), Inches(0.95))

# Cedars ref
rect(s, Inches(7.05), Inches(5.65), Inches(5.6), Inches(0.5),
     fill=NAVY_BG, border=NAVY, bw=Pt(1))
label(s, "Cedars reference: BERT=79.6% (T)  |  LLaMA zero-shot: T=54%, N=84%, M=91%",
      Inches(7.2), Inches(5.72), Inches(5.3), Inches(0.35),
      size=11, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE B — WHAT IS TNM STAGING?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)

rect(s, Inches(0), Inches(0), Inches(13.33), Inches(1.1), fill=NAVY)
label(s, "What is TNM Staging?",
      Inches(0.4), Inches(0.12), Inches(10), Inches(0.6),
      size=26, bold=True, color=WHITE)
label(s, "A standard system used by pathologists worldwide to measure how serious a cancer is",
      Inches(0.4), Inches(0.72), Inches(12), Inches(0.32),
      size=13, color=RGBColor(0xa8,0xc8,0xe8))

# T box
rect(s, Inches(0.4), Inches(1.3), Inches(3.9), Inches(3.8),
     fill=BLUE_BG, border=NAVY, bw=Pt(2))
label(s, "T — Tumor Size",
      Inches(0.55), Inches(1.4), Inches(3.6), Inches(0.45),
      size=18, bold=True, color=NAVY)
label(s, "How big is the primary tumor?\nHow deep has it grown?",
      Inches(0.55), Inches(1.88), Inches(3.6), Inches(0.55),
      size=12, italic=True, color=MED_GRAY)
rect(s, Inches(0.5), Inches(2.5), Inches(3.7), Pt(1.5), fill=NAVY)
tstages = [
    ("T1", "Small tumor, still contained"),
    ("T2", "Larger, growing into nearby tissue"),
    ("T3", "Large, grown into surrounding area"),
    ("T4", "Very large, invaded nearby organs"),
]
for i, (code, desc) in enumerate(tstages):
    t = Inches(2.62 + i * 0.6)
    rect(s, Inches(0.55), t, Inches(0.45), Inches(0.48), fill=NAVY_BG)
    label(s, code, Inches(0.55), t + Inches(0.08),
          Inches(0.45), Inches(0.35), size=12, bold=True,
          color=NAVY, align=PP_ALIGN.CENTER)
    label(s, desc, Inches(1.1), t + Inches(0.1),
          Inches(3.1), Inches(0.35), size=12, color=DARK_TEXT)

# N box
rect(s, Inches(4.7), Inches(1.3), Inches(3.9), Inches(3.8),
     fill=PURPLE_BG, border=PURPLE, bw=Pt(2))
label(s, "N — Lymph Nodes",
      Inches(4.85), Inches(1.4), Inches(3.6), Inches(0.45),
      size=18, bold=True, color=PURPLE)
label(s, "Has cancer spread to nearby\nlymph nodes?",
      Inches(4.85), Inches(1.88), Inches(3.6), Inches(0.55),
      size=12, italic=True, color=MED_GRAY)
rect(s, Inches(4.8), Inches(2.5), Inches(3.7), Pt(1.5), fill=PURPLE)
nstages = [
    ("N0", "No lymph nodes affected"),
    ("N1", "Cancer in 1-3 nearby nodes"),
    ("N2", "Cancer in 4-9 nearby nodes"),
    ("N3", "Cancer in many lymph nodes"),
]
for i, (code, desc) in enumerate(nstages):
    t = Inches(2.62 + i * 0.6)
    rect(s, Inches(4.85), t, Inches(0.45), Inches(0.48), fill=PURPLE_BG)
    label(s, code, Inches(4.85), t + Inches(0.08),
          Inches(0.45), Inches(0.35), size=12, bold=True,
          color=PURPLE, align=PP_ALIGN.CENTER)
    label(s, desc, Inches(5.4), t + Inches(0.1),
          Inches(3.1), Inches(0.35), size=12, color=DARK_TEXT)

# M box
rect(s, Inches(9.0), Inches(1.3), Inches(3.9), Inches(3.8),
     fill=RED_BG, border=RED, bw=Pt(2))
label(s, "M — Metastasis",
      Inches(9.15), Inches(1.4), Inches(3.6), Inches(0.45),
      size=18, bold=True, color=RED)
label(s, "Has cancer spread to distant\norgans (lungs, liver, bones)?",
      Inches(9.15), Inches(1.88), Inches(3.6), Inches(0.55),
      size=12, italic=True, color=MED_GRAY)
rect(s, Inches(9.1), Inches(2.5), Inches(3.7), Pt(1.5), fill=RED)
mstages = [
    ("M0", "No spread to distant organs"),
    ("M1", "Has spread to distant organs"),
    ("", ""),
    ("", "M1 = most serious stage"),
]
for i, (code, desc) in enumerate(mstages):
    if not code and not desc:
        continue
    t = Inches(2.62 + i * 0.6)
    if code:
        rect(s, Inches(9.15), t, Inches(0.45), Inches(0.48), fill=RED_BG)
        label(s, code, Inches(9.15), t + Inches(0.08),
              Inches(0.45), Inches(0.35), size=12, bold=True,
              color=RED, align=PP_ALIGN.CENTER)
    label(s, desc, Inches(9.7), t + Inches(0.1),
          Inches(3.1), Inches(0.35), size=12,
          color=RED if not code else DARK_TEXT, italic=(not code))

# Example
rect(s, Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.85),
     fill=GREEN_BG, border=GREEN, bw=Pt(1.5))
label(s, "Real Example:",
      Inches(0.6), Inches(5.4), Inches(1.5), Inches(0.35),
      size=14, bold=True, color=GREEN)
label(s, "BRCA patient — T2  N1  M0   means:",
      Inches(2.1), Inches(5.4), Inches(4.0), Inches(0.35),
      size=14, bold=True, color=DARK_TEXT)
label(s, "Medium tumor (T2)  +  some lymph nodes affected (N1)  +  NOT spread elsewhere (M0)  =  Stage 2 Breast Cancer",
      Inches(0.6), Inches(5.78), Inches(12.0), Inches(0.3),
      size=13, color=DARK_TEXT)

# How labels were made
rect(s, Inches(0.4), Inches(6.32), Inches(12.5), Inches(0.95),
     fill=LIGHT_GRAY, border=BORDER, bw=Pt(1))
label(s, "How TNM labels were created:",
      Inches(0.6), Inches(6.4), Inches(3.2), Inches(0.35),
      size=13, bold=True, color=NAVY)
label(s, "Real pathologists read each report and manually assigned T, N, M values. "
         "These are stored in the TNM_Stage CSV files. "
         "Our model's job: read the same text and predict what the pathologist would write — automatically.",
      Inches(3.8), Inches(6.4), Inches(9.0), Inches(0.75),
      size=12, color=DARK_TEXT)


# ── SAVE ──────────────────────────────────────────────────────────────────────
OUT = (r"c:\Users\kings\OneDrive\Desktop\CSUDH_Kingslee"
       r"\CSUDH-Spring-2026\CSC-590-MastersProject"
       r"\MastersProject\TNM_Slides.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
print("2 slides created:")
print("  Slide 1 — Task 1 vs Task 2 side by side comparison")
print("  Slide 2 — What is TNM staging? (T, N, M explained)")
