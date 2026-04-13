"""
Creates 2 slides:
  Slide 1 — BoW vs TF-IDF vs LR explained visually
  Slide 2 — What reference did vs What we are doing
Run: python create_method_slides.py
Output: Method_Slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
DARK_TEXT  = RGBColor(0x1e, 0x29, 0x3b)
MED_GRAY   = RGBColor(0x47, 0x55, 0x69)
LIGHT_GRAY = RGBColor(0xf1, 0xf5, 0xf9)
BORDER     = RGBColor(0xcb, 0xd5, 0xe1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x16, 0xa3, 0x4a)
GREEN_BG   = RGBColor(0xdc, 0xfc, 0xe7)
ORANGE     = RGBColor(0xea, 0x58, 0x0c)
ORANGE_BG  = RGBColor(0xff, 0xed, 0xd5)
BLUE_BG    = RGBColor(0xef, 0xf6, 0xff)
NAVY_BG    = RGBColor(0xe0, 0xe7, 0xf1)
PURPLE     = RGBColor(0x7c, 0x3a, 0xed)
PURPLE_BG  = RGBColor(0xed, 0xe9, 0xfe)
RED        = RGBColor(0xdc, 0x26, 0x26)
RED_BG     = RGBColor(0xfe, 0xe2, 0xe2)
YELLOW_BG  = RGBColor(0xff, 0xf9, 0xc4)
YELLOW     = RGBColor(0xca, 0x8a, 0x04)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def label(slide, text, l, t, w, h, size=13, bold=False,
          color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def multiline(slide, lines, l, t, w, h, ds=13, dc=DARK_TEXT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, cfg) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = cfg.get("align", PP_ALIGN.LEFT)
        p.space_after  = Pt(cfg.get("sa", 3))
        p.space_before = Pt(cfg.get("sb", 0))
        r = p.add_run()
        r.text = text
        r.font.size    = Pt(cfg.get("size", ds))
        r.font.bold    = cfg.get("bold", False)
        r.font.italic  = cfg.get("italic", False)
        r.font.color.rgb = cfg.get("color", dc)

def top_bar(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), fill=NAVY)
    label(slide, title,
          Inches(0.4), Inches(0.12), Inches(12), Inches(0.6),
          size=26, bold=True, color=WHITE)
    if subtitle:
        label(slide, subtitle,
              Inches(0.4), Inches(0.72), Inches(12), Inches(0.32),
              size=13, color=RGBColor(0xa8,0xc8,0xe8))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — BoW vs TF-IDF vs Logistic Regression explained
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "BoW vs TF-IDF vs Logistic Regression — Explained",
        "How the reference pipeline works and how we improved it")

# ── SECTION: BoW ──
rect(s, Inches(0.35), Inches(1.25), Inches(3.9), Inches(5.85),
     fill=ORANGE_BG, border=ORANGE, bw=Pt(1.5))

label(s, "Bag of Words  (BoW)",
      Inches(0.5), Inches(1.35), Inches(3.6), Inches(0.42),
      size=15, bold=True, color=ORANGE)
label(s, "Used by Cedars-Sinai reference",
      Inches(0.5), Inches(1.78), Inches(3.6), Inches(0.28),
      size=11, italic=True, color=MED_GRAY)
rect(s, Inches(0.45), Inches(2.1), Inches(3.7), Pt(1.5), fill=ORANGE)

# BoW example
rect(s, Inches(0.5), Inches(2.22), Inches(3.6), Inches(0.6),
     fill=WHITE, border=BORDER)
label(s, '"The tumor is invasive. The tumor is large."',
      Inches(0.6), Inches(2.3), Inches(3.4), Inches(0.38),
      size=10, italic=True, color=DARK_TEXT)

label(s, "Counts how many times each word appears:",
      Inches(0.5), Inches(2.95), Inches(3.6), Inches(0.32),
      size=11, color=DARK_TEXT)

# Word count table
words = [("the", "2"), ("tumor", "2"), ("is", "2"), ("invasive", "1"), ("large", "1")]
for i, (word, count) in enumerate(words):
    l = Inches(0.5 + i * 0.72)
    rect(s, l, Inches(3.35), Inches(0.65), Inches(0.35),
         fill=NAVY_BG, border=BORDER)
    label(s, word, l, Inches(3.37), Inches(0.65), Inches(0.28),
          size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s, l, Inches(3.72), Inches(0.65), Inches(0.32),
         fill=WHITE, border=BORDER)
    label(s, count, l, Inches(3.74), Inches(0.65), Inches(0.28),
          size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

label(s, "Problem: 'the' and 'tumor' both count = 2\nbut 'the' tells us NOTHING about cancer type!",
      Inches(0.5), Inches(4.18), Inches(3.6), Inches(0.6),
      size=11, color=RED, italic=True)

multiline(s, [
    ("Weaknesses of BoW:", {"bold": True, "size": 12, "color": ORANGE, "sa": 4}),
    ("  Common words (the, is, a) inflate counts", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  Rare but important words get buried", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  All words treated equally — no importance ranking", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  Result: 95.31% on 32 classes", {"size": 12, "color": ORANGE, "bold": True, "sa": 4}),
], Inches(0.5), Inches(4.88), Inches(3.6), Inches(1.9), ds=11)


# ── SECTION: TF-IDF ──
rect(s, Inches(4.55), Inches(1.25), Inches(3.9), Inches(5.85),
     fill=BLUE_BG, border=NAVY, bw=Pt(1.5))

label(s, "TF-IDF",
      Inches(4.7), Inches(1.35), Inches(3.6), Inches(0.42),
      size=15, bold=True, color=NAVY)
label(s, "Term Frequency - Inverse Document Frequency",
      Inches(4.7), Inches(1.78), Inches(3.6), Inches(0.28),
      size=11, italic=True, color=MED_GRAY)
rect(s, Inches(4.6), Inches(2.1), Inches(3.7), Pt(1.5), fill=NAVY)

label(s, "TF  (Term Frequency):",
      Inches(4.7), Inches(2.22), Inches(3.6), Inches(0.32),
      size=12, bold=True, color=NAVY)
label(s, "How often does this word appear in THIS report?",
      Inches(4.7), Inches(2.55), Inches(3.6), Inches(0.38),
      size=11, color=DARK_TEXT)

label(s, "IDF  (Inverse Document Frequency):",
      Inches(4.7), Inches(3.02), Inches(3.6), Inches(0.32),
      size=12, bold=True, color=NAVY)
label(s, "How RARE is this word across ALL reports?\nRare = more important. Common = less important.",
      Inches(4.7), Inches(3.35), Inches(3.6), Inches(0.52),
      size=11, color=DARK_TEXT)

label(s, "TF-IDF Score  =  TF  x  IDF",
      Inches(4.7), Inches(4.0), Inches(3.6), Inches(0.38),
      size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Score comparison
scores = [
    ("the",        "HIGH TF", "LOW IDF",  "LOW score",  RED_BG,   RED),
    ("adenocarcinoma", "MED TF", "HIGH IDF", "HIGH score", GREEN_BG, GREEN),
]
for i, (word, tf, idf, result, bg, col) in enumerate(scores):
    t = Inches(4.52 + i * 0.72)
    rect(s, Inches(4.7), t, Inches(3.6), Inches(0.6), fill=bg, border=col, bw=Pt(1))
    label(s, f'"{word}"  {tf} x {idf} = {result}',
          Inches(4.8), t + Inches(0.12), Inches(3.4), Inches(0.35),
          size=11, bold=True, color=col)

multiline(s, [
    ("Advantages over BoW:", {"bold": True, "size": 12, "color": NAVY, "sa": 4}),
    ("  Common words auto-downweighted", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  Medical terms like 'adenocarcinoma' ranked HIGH", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  15,000 features + bigrams (word pairs)", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  Result: 96.36% on 35 classes (beats ref)", {"size": 12, "color": GREEN, "bold": True, "sa": 4}),
], Inches(4.7), Inches(5.55), Inches(3.6), Inches(1.4), ds=11)


# ── SECTION: Logistic Regression ──
rect(s, Inches(8.75), Inches(1.25), Inches(4.2), Inches(5.85),
     fill=PURPLE_BG, border=PURPLE, bw=Pt(1.5))

label(s, "Logistic Regression  (LR)",
      Inches(8.9), Inches(1.35), Inches(3.9), Inches(0.42),
      size=15, bold=True, color=PURPLE)
label(s, "The classifier — makes the final decision",
      Inches(8.9), Inches(1.78), Inches(3.9), Inches(0.28),
      size=11, italic=True, color=MED_GRAY)
rect(s, Inches(8.8), Inches(2.1), Inches(4.0), Pt(1.5), fill=PURPLE)

label(s, "What it does:",
      Inches(8.9), Inches(2.22), Inches(3.9), Inches(0.32),
      size=12, bold=True, color=PURPLE)
multiline(s, [
    ("Takes the TF-IDF numbers as input.", {"size": 11, "color": DARK_TEXT, "sa": 5}),
    ("Learns which words strongly predict each cancer type.", {"size": 11, "color": DARK_TEXT, "sa": 5}),
    ("Outputs a probability for all 35 classes.", {"size": 11, "color": DARK_TEXT, "sa": 5}),
    ("Picks the class with the highest probability.", {"size": 11, "color": DARK_TEXT, "sa": 8}),
], Inches(8.9), Inches(2.58), Inches(3.9), Inches(1.4), ds=11)

# Pipeline visual
label(s, "Full pipeline:",
      Inches(8.9), Inches(4.1), Inches(3.9), Inches(0.32),
      size=12, bold=True, color=PURPLE)
steps = [
    ("Report Text", WHITE, BORDER),
    ("TF-IDF scores", NAVY_BG, NAVY),
    ("LR predicts", PURPLE_BG, PURPLE),
    ("Cancer Type", GREEN_BG, GREEN),
]
for i, (step, bg, col) in enumerate(steps):
    t = Inches(4.5 + i * 0.65)
    rect(s, Inches(8.9), t, Inches(3.9), Inches(0.52),
         fill=bg, border=col, bw=Pt(1))
    label(s, step, Inches(8.9), t + Inches(0.1),
          Inches(3.9), Inches(0.35), size=12,
          color=col if col != BORDER else DARK_TEXT,
          align=PP_ALIGN.CENTER, bold=True)
    if i < 3:
        label(s, "v", Inches(10.55), t + Inches(0.52),
              Inches(0.3), Inches(0.2),
              size=9, color=MED_GRAY, align=PP_ALIGN.CENTER)

multiline(s, [
    ("Why LR works well for text:", {"bold": True, "size": 12, "color": PURPLE, "sa": 4}),
    ("  Fast to train — seconds, not hours", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  Explainable — you can see which words", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  drove the prediction", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  No GPU needed", {"size": 11, "color": DARK_TEXT, "sa": 4}),
    ("  Baseline to beat with BERT and RAG", {"size": 11, "color": DARK_TEXT, "sa": 4}),
], Inches(8.9), Inches(7.0) - Inches(1.6), Inches(3.9), Inches(1.5), ds=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Reference vs Our Work
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg_white(s)
top_bar(s, "What the Reference Did vs What We Are Building",
        "Cedars-Sinai provided a baseline — MEDREx goes significantly further")

# Column headers
rect(s, Inches(0.35), Inches(1.2), Inches(6.1), Inches(0.45), fill=ORANGE_BG, border=ORANGE, bw=Pt(1))
label(s, "Cedars-Sinai Reference  (what was given to us)",
      Inches(0.5), Inches(1.28), Inches(5.8), Inches(0.32),
      size=14, bold=True, color=ORANGE)

rect(s, Inches(6.9), Inches(1.2), Inches(6.1), Inches(0.45), fill=BLUE_BG, border=NAVY, bw=Pt(1))
label(s, "MEDREx  (what we are building)",
      Inches(7.05), Inches(1.28), Inches(5.8), Inches(0.32),
      size=14, bold=True, color=NAVY)

# Divider
rect(s, Inches(6.65), Inches(1.2), Pt(2), Inches(6.1), fill=BORDER)

# Task 1 section
rect(s, Inches(0.35), Inches(1.78), Inches(12.65), Inches(0.35),
     fill=LIGHT_GRAY, border=BORDER, bw=Pt(0.5))
label(s, "TASK 1 — Cancer Type Classification",
      Inches(0.5), Inches(1.83), Inches(6), Inches(0.25),
      size=12, bold=True, color=NAVY)

rows_t1 = [
    # (ref_text, ref_col, our_text, our_col)
    ("BoW (raw word counts only)",
     MED_GRAY,
     "TF-IDF (weighted — medical terms ranked higher)",
     GREEN),
    ("Logistic Regression only",
     MED_GRAY,
     "LR + SVM + BERT + Embedding Similarity + RAG+LLM",
     GREEN),
    ("32 cancer types covered",
     MED_GRAY,
     "35 cancer types — 3 more classes",
     GREEN),
    ("95.31% accuracy",
     ORANGE,
     "96.36% already — and still improving",
     GREEN),
    ("No vector database",
     MED_GRAY,
     "ChromaDB — 9,523 reports indexed, semantic search live",
     GREEN),
    ("No RAG pipeline",
     MED_GRAY,
     "RAG + LLM pipeline planned — our novel contribution",
     NAVY),
]
for i, (ref, rc, ours, oc) in enumerate(rows_t1):
    t = Inches(2.28 + i * 0.57)
    bg_ref = RED_BG if rc == MED_GRAY else ORANGE_BG
    bg_our = GREEN_BG if oc == GREEN else BLUE_BG
    rect(s, Inches(0.35), t, Inches(6.1), Inches(0.5),
         fill=bg_ref, border=BORDER, bw=Pt(0.5))
    label(s, ref, Inches(0.5), t + Inches(0.1),
          Inches(5.8), Inches(0.35), size=12, color=rc)
    rect(s, Inches(6.9), t, Inches(6.1), Inches(0.5),
         fill=bg_our, border=BORDER, bw=Pt(0.5))
    label(s, ours, Inches(7.05), t + Inches(0.1),
          Inches(5.8), Inches(0.35), size=12, color=oc)

# Task 2 section
rect(s, Inches(0.35), Inches(5.76), Inches(12.65), Inches(0.35),
     fill=LIGHT_GRAY, border=BORDER, bw=Pt(0.5))
label(s, "TASK 2 — TNM Staging",
      Inches(0.5), Inches(5.81), Inches(6), Inches(0.25),
      size=12, bold=True, color=NAVY)

rows_t2 = [
    ("BoW + LR baseline",
     MED_GRAY,
     "TF-IDF + LR baseline  (to build next)",
     NAVY),
    ("BERT for T-stage → 79.6%   N/M also done",
     ORANGE,
     "Reproduce + improve with our pipeline",
     NAVY),
    ("LLaMA zero-shot → T=54%, N=84%, M=91%",
     ORANGE,
     "Reproduce + add RAG context for better accuracy",
     NAVY),
]
for i, (ref, rc, ours, oc) in enumerate(rows_t2):
    t = Inches(6.22 + i * 0.57)
    rect(s, Inches(0.35), t, Inches(6.1), Inches(0.5),
         fill=ORANGE_BG, border=BORDER, bw=Pt(0.5))
    label(s, ref, Inches(0.5), t + Inches(0.1),
          Inches(5.8), Inches(0.35), size=12, color=rc)
    rect(s, Inches(6.9), t, Inches(6.1), Inches(0.5),
         fill=BLUE_BG, border=BORDER, bw=Pt(0.5))
    label(s, ours, Inches(7.05), t + Inches(0.1),
          Inches(5.8), Inches(0.35), size=12, color=oc)


# ── SAVE ──────────────────────────────────────────────────────────────────────
OUT = (r"c:\Users\kings\OneDrive\Desktop\CSUDH_Kingslee"
       r"\CSUDH-Spring-2026\CSC-590-MastersProject"
       r"\MastersProject\Method_Slides.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
print("2 slides:")
print("  Slide 1 - BoW vs TF-IDF vs Logistic Regression explained")
print("  Slide 2 - Reference vs MEDREx comparison (Task 1 + Task 2)")
